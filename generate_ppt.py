"""
Usage instructions:
1. Install dependencies: ``pip install -r requirements.txt``
2. Paste the full essay text into ``ESSAY_TEXT`` below.
3. Run the script: ``python generate_ppt.py``

This script will create a title slide from the essay header, generate one
content slide per section with concise bullet points, and add an image
placeholder on each section slide.
"""
from __future__ import annotations

import re
from typing import Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ESSAY_TEXT = """erspectives on Idiopathic Short Stature: An Academic Essay for Medical Trainees 
Prepared for: Mansoura University – Faculty of Medicine
Student: [Your Name Here]
Course: [Course Name]
Date: [Month, Year]

1. Introduction

Idiopathic short stature (ISS) is usually defined as height more than two standard deviations below the mean for age and sex in a child who has no identifiable chronic disease, endocrine disorder, or nutritional deficiency. In practice, ISS is a working label that is applied after the main known causes of growth failure have been excluded, while the precise biological mechanism remains uncertain. This definition underpins current international guidelines and reminds clinicians that ISS is not a single disease but a heterogeneous group of conditions that share a similar growth pattern.
Reference: Cohen P, Rogol AD, Deal CL, et al. Consensus statement on the diagnosis and treatment of children with idiopathic short stature. J Clin Endocrinol Metab. 2008;93(11):4210–4217.

In pediatric endocrine clinics, children with suspected ISS form a large proportion of referrals, yet many short children in the community are never evaluated. Whether a child is investigated often depends on parental concern, observations from schools or primary care, and access to specialist services. This means that ISS is influenced by healthcare systems and social factors as well as biology.

For medical trainees, ISS is a useful model of how to handle diagnostic uncertainty. Good practice depends on careful history-taking, repeated growth measurements, and clear explanations to families about what is known, what is not yet known, and why follow-up is necessary. The diagnosis also forces clinicians to think about psychosocial impact and financial cost, not only about centimeters gained on a growth chart.

2. Normal Growth and Physiology

Normal linear growth reflects the interaction of genetics, nutrition, hormones, and the local biology of the growth plate. The growth hormone (GH)–insulin-like growth factor-1 (IGF-1) axis is central: GH secreted in pulses from the anterior pituitary stimulates hepatic and local IGF-1 production, and IGF-1 then acts on growth plate chondrocytes to support proliferation and hypertrophy. Understanding this axis helps clinicians interpret IGF-1 levels, GH stimulation tests, and the likely response to treatment in children with short stature.
Reference: Rogol AD, Hayden GF. Etiologies and early diagnosis of short stature and growth failure in children and adolescents. Pediatr Clin North Am. 2015;62(4):1019–1036.

Growth velocity is not constant. Infancy is marked by rapid growth, followed by a slower and steady phase in mid-childhood and a pubertal growth spurt driven by sex steroids and rising GH secretion. Interruption at any stage, even by a seemingly small amount, can reduce final adult height. Adequate calories, protein, and micronutrients such as zinc and vitamin D are needed throughout childhood, and chronic illness can divert energy away from growth toward immune or metabolic demands.

Sleep and daily routines also influence growth. GH secretion peaks during deep sleep, so poor sleep quality, irregular schedules, or chronic stress can reduce effective exposure of the growth plate to GH. Lifestyle counseling about regular sleep, balanced nutrition, and management of chronic disease is therefore part of growth care, even when no specific endocrine defect is found.

Anthropometry itself is a core physiological tool. Accurate, repeated measurements with a calibrated stadiometer, plotted on appropriate growth charts and compared with mid-parental height, are essential. Many children who appear short at a single visit are later found to be following a stable familial pattern once several points over time are available, avoiding unnecessary investigations.

3. Pathophysiology and Genetics of Idiopathic Short Stature

Historically, ISS was seen mainly as a diagnosis of exclusion. However, genetic and molecular studies show that many children labeled with ISS have subtle defects affecting growth regulation. Some have variants in single genes that influence growth plate architecture, GH secretion, or IGF-1 signaling. Others have a complex polygenic background, in which many common variants each contribute a small reduction in height. These discoveries blur the boundary between “idiopathic” and “monogenic” short stature.
Reference: Dauber A, Rosenfeld RG, Hirschhorn JN. Genetic evaluation of short stature. Nat Rev Endocrinol. 2014;10(10):582–593.

Variants in genes such as ACAN, NPR2, or SHOX can cause proportionate short stature without striking skeletal deformities, so the child may look “normal but small”. Without genetic testing, such cases are often grouped under ISS. Subtle resistance to GH or IGF-1, due to changes in receptors or downstream signaling pathways, can also lead to poor growth despite apparently normal GH stimulation test results, because signaling inside the cell is impaired even when hormone levels appear adequate.

Epigenetic mechanisms and polygenic scores add further complexity. Environmental factors such as intrauterine growth restriction, chronic undernutrition, and psychosocial stress can modify gene expression without altering the DNA sequence. This may help explain why two children with similar genetic backgrounds and similar height at one time point can follow different growth trajectories over time.

Clinically, it is useful to think of ISS as a spectrum of overlapping mechanisms rather than a single entity. This perspective supports selective use of genetic testing in children with strong family history, unusual body proportions, dysmorphic features, or poor response to standard therapy, while acknowledging that many children will still have no clear molecular diagnosis despite careful evaluation.

4. Diagnostic Evaluation and Workup

The evaluation of suspected ISS begins with careful history, physical examination, and auxologic assessment before extensive blood testing is considered. Important elements include birth weight and length, early feeding and illnesses, chronic symptoms such as gastrointestinal or respiratory problems, medication use, and a detailed family history of height and pubertal timing. Serial height measurements allow calculation of growth velocity and comparison with mid-parental target height, which helps estimate how far the child is from genetic expectations.
Reference: Grimberg A, DiVall SA, Polychronakos C, et al. Guidelines for growth hormone and insulin-like growth factor-I treatment in children and adolescents. Horm Res Paediatr. 2016;86(6):361–397.

Baseline laboratory tests often include a complete blood count, renal and liver function, thyroid function, celiac serology, and IGF-1 levels. More specialized tests, such as GH stimulation testing or adrenal studies, are generally reserved for children with poor growth velocity, abnormal screening tests, or other clinical signs suggestive of specific endocrine disorders. Bone age radiographs help distinguish constitutional delay of growth and puberty from more concerning growth failure and provide an estimate of remaining growth potential.

Body proportions, facial features, and pubertal staging can give clues to syndromic or skeletal diagnoses that would not fit ISS. In selected cases, targeted or panel-based genetic testing is appropriate, especially when there is a strong familial pattern, disproportion, or dysmorphism. Throughout the process, it is important to explain each step, the likely yield of tests, and reasons for watchful waiting, so that families understand why investigations are paced over time.

Regular follow-up, typically every four to six months, is essential because growth patterns can change and puberty may begin. ISS is therefore not a fixed or permanent label but a provisional diagnosis that should be revisited whenever new clinical information arises.

5. Psychosocial and Quality-of-Life Aspects

Short stature can influence self-image and social experience, even when physical health is otherwise normal. Some children with marked short stature report more teasing, social exclusion, or self-consciousness about appearance, although others adapt well and have no major difficulties. The impact depends on school culture, family support, and individual temperament, as well as height itself.
Reference: Wheeler PG, Bresnahan K, Shephard BA, Lau J, Balk EM. Short stature and functional impairment: a systematic review. J Pediatr. 2004;144(4):364–370.

Parents and teachers may unintentionally reinforce differences by overprotecting the child or lowering expectations for independence. This can limit opportunities to build confidence and practical skills. Education that separates stature from competence, gives age-appropriate responsibilities, and avoids labeling the child as fragile can reduce this effect. Simple adjustments at school, such as seating position, access to equipment, or choice of sports, often make daily life easier.

Psychological support can be valuable when there is clear distress, bullying, or social withdrawal. Counseling helps children reframe height as one feature of identity rather than the main defining characteristic and encourages them to recognize and build on personal strengths. Meeting peers with similar experiences, whether locally or through support groups, can reduce isolation and normalize concerns.

During endocrine visits, asking directly but sensitively about school, friendships, and mood ensures that psychosocial issues are not ignored in favor of laboratory values and imaging. A brief screening conversation may reveal problems that are more important to quality of life than any potential height gain.

6. Therapeutic Options and Monitoring

Recombinant human growth hormone (rhGH) is the main pharmacologic treatment offered to many children with ISS in settings where it is available. On average, it can increase adult height by several centimeters when started in mid-childhood and continued until near final height, but the response is highly variable. Decisions to treat should therefore be based on careful discussion of expected benefit, the burden of daily injections, and the uncertainty in individual outcomes.
Reference: Ranke MB. Growth hormone therapy in idiopathic short stature. Best Pract Res Clin Endocrinol Metab. 2015;29(3):353–366.

rhGH treatment requires daily subcutaneous injections, regular clinic visits, and considerable financial cost. Monitoring includes growth velocity, IGF-1 levels, and periodic bone age, as well as screening for adverse effects such as intracranial hypertension, slipped capital femoral epiphysis, and disturbances of glucose metabolism. When growth response is poor, clinicians should review adherence, injection technique, and the original diagnosis before increasing dose or extending treatment duration.

Adjunctive options, such as gonadotropin-releasing hormone analogues in early or rapidly progressing puberty, or aromatase inhibitors in selected boys, aim to prolong the period of growth by slowing bone maturation. These strategies carry their own risks and uncertainties, so they are usually reserved for specific scenarios and should be discussed thoroughly with families, with clear documentation of reasons for use.

Non-pharmacologic measures remain important regardless of rhGH use. Optimizing nutrition, managing chronic illnesses, encouraging regular sleep, and supporting mental health help the child make the best use of any remaining growth potential. Transition planning for older adolescents should include when to stop treatment, how to monitor weight and metabolic health in adult care, and how to cope with final adult height in a realistic and supportive way.

7. Ethics, Health Economics, and Policy

The use of rhGH in ISS raises ethical and economic questions. Short stature lies at the lower end of a normal distribution and is not a life-threatening condition, while height gains are usually modest. This has led to debate about whether treating ISS sometimes medicalizes normal variation, and whether high-cost therapy is justified in all cases. Clinicians must balance potential psychosocial benefits against medical risks, financial costs, and the message society sends about body size.
Reference: Allen DB, Fost N. Growth hormone therapy for short stature: panacea or peril? N Engl J Med. 2004;351(16):1610–1614.

Health-economic analyses suggest that the cost per centimeter gained in ISS is high compared with many other pediatric interventions. In healthcare systems with finite resources, widespread use of rhGH for ISS may compete with other priorities. For this reason, many countries and insurers apply eligibility criteria based on height standard deviation scores, growth velocity, and predicted adult height before approving treatment. Clear documentation and shared decision-making help ensure that these criteria are applied fairly and are understood by families.

Equity of access is another concern. Families with more knowledge, time, or resources may be more likely to seek evaluation and push for treatment, while others with similar clinical need may never reach a specialist. Policies that support early growth monitoring in primary care, clear referral pathways, and accessible patient information can reduce these differences, making decisions less dependent on background and advocacy skills.

Ethically sound care also requires honest communication about uncertainty. Long-term psychosocial gains are not guaranteed, and some children may adapt well without treatment. Clinicians should explain alternative approaches, such as psychological support, and emphasize that choosing not to treat with rhGH is a valid option when benefits are unclear or treatment burdens are high.

8. Future Directions and Research Frontiers

New genetic and genomic technologies are steadily changing how short stature is classified. Next-generation sequencing, copy-number analysis, and functional studies of genes active in chondrocytes have identified several novel causes of short stature, and this process is ongoing. Over time, these discoveries are likely to divide ISS into smaller, better defined diagnostic groups with different natural histories and treatment responses.
Reference: Wit JM, Oostdijk W, Losekoot M, van Duyvenvoorde HA, Ruivenkamp CA, Kant SG. Novel genetic causes of short stature. Eur J Endocrinol. 2016;174(6):R145–R173.

Pharmaceutical research is exploring long-acting GH preparations, new IGF-1 formulations, and drugs that act directly on growth plate signaling pathways. Long-acting GH aims to reduce injection frequency and possibly improve adherence, but long-term safety and real-world effectiveness still need careful evaluation. In the future, therapies may be tailored more closely to specific genetic or signaling defects rather than using a single regimen for all children with ISS.

Digital health and large clinical registries are also likely to shape ISS care. Mobile applications and connected devices can help families track injections, sleep, physical activity, and growth, generating detailed data that may be linked with clinic records. International registries that collect information on height outcomes, side effects, and quality of life in routine practice will be important for refining guidelines and for understanding which children benefit most from treatment.

Research into psychosocial outcomes and health economics will need to keep pace with biological advances. Future studies that combine endocrine, psychological, and educational perspectives can clarify which interventions truly improve daily functioning and wellbeing, helping clinicians and policy makers make informed decisions about how to use new technologies.

9. Conclusion

Idiopathic short stature illustrates the intersection of biological complexity, diagnostic uncertainty, and patient-centered care. Evaluation requires careful auxology, thoughtful use of laboratory and imaging tests, and willingness to revisit the diagnosis as new information appears. Management extends beyond the prescription of rhGH to include attention to psychosocial wellbeing, education of families and schools, and realistic discussion of what treatment can and cannot achieve.

Because responses to therapy are heterogeneous and long-term psychosocial outcomes are variable, individualized care plans are essential. These plans should align therapeutic intensity with the child’s and family’s priorities, while remaining mindful of safety, cost, and equity. Shared decision-making provides a framework for these conversations and helps maintain trust over years of follow-up.

Ongoing advances in genetics, pharmacology, and digital health will continue to change what is possible in ISS, but the core tasks for clinicians will remain the same: listening carefully, measuring accurately, explaining honestly, and advocating for fair access to appropriate care. For medical trainees, ISS offers a practical example of how scientific knowledge, ethical reasoning, and empathy must come together in everyday clinical practice.
"""


SectionData = Dict[str, List[str]]
ParsedEssay = Dict[str, object]


def parse_essay(essay_text: str) -> ParsedEssay:
    """Parse the essay text into a structured dictionary."""
    text = essay_text.strip()
    if not text:
        raise ValueError("Essay text is empty. Please paste the full essay into ESSAY_TEXT.")

    lines = text.splitlines()
    header_lines: List[str] = []
    index = 0
    while index < len(lines) and len(header_lines) < 5:
        current = lines[index].strip()
        if current:
            header_lines.append(current)
        index += 1

    if len(header_lines) < 5:
        raise ValueError("Expected at least five header lines (title, prepared for, student, course, date).")

    header = {
        "title": header_lines[0],
        "prepared_for": header_lines[1],
        "student": header_lines[2],
        "course": header_lines[3],
        "date": header_lines[4],
    }

    remaining_text = "\n".join(lines[index:]).strip()
    if not remaining_text:
        raise ValueError("No section content found after the header block.")

    section_pattern = re.compile(r"^\d+\.\s.+$", re.MULTILINE)
    matches = list(section_pattern.finditer(remaining_text))
    if not matches:
        raise ValueError("No numbered sections found in the essay text.")

    sections: List[Dict[str, object]] = []
    for idx, match in enumerate(matches):
        heading = match.group().strip()
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(remaining_text)
        body = remaining_text[start:end].strip()
        paragraphs = [para.strip() for para in re.split(r"\n\s*\n", body) if para.strip()]
        sections.append({"heading": heading, "paragraphs": paragraphs})

    return {"title_block": header, "sections": sections}


def _extract_bullet_points(paragraphs: List[str], limit: int = 4) -> List[str]:
    """Generate bullet points from the given paragraphs."""
    bullets: List[str] = []
    for paragraph in paragraphs:
        if len(bullets) >= limit:
            break

        main_text = paragraph.split("Reference:")[0].strip()
        if not main_text:
            continue

        period_index = main_text.find(".")
        if period_index != -1:
            first_sentence = main_text[: period_index + 1].strip()
        else:
            first_sentence = main_text.strip()

        if not first_sentence:
            continue

        if len(first_sentence) > 150:
            truncated = first_sentence[:147].rstrip()
            first_sentence = f"{truncated}..."

        bullets.append(first_sentence)

    return bullets


def _configure_title_shape(shape, text: str, font_size: int = 36) -> None:
    """Apply consistent styling to a title shape."""
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)


def _configure_body_placeholder(shape, bullets: List[str]) -> None:
    """Populate the content placeholder with bullet points."""
    text_frame = shape.text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.name = "Calibri"
        paragraph.font.size = Pt(22)
        paragraph.level = 0


def _add_image_placeholder(slide) -> None:
    """Add a light rectangle placeholder for images on the slide."""
    left = Inches(6.5)
    top = Inches(1.5)
    width = Inches(3)
    height = Inches(4.5)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)
    line = shape.line
    line.color.rgb = RGBColor(200, 200, 200)
    text_frame = shape.text_frame
    text_frame.text = "Insert relevant image here\n(e.g., growth chart, physiology,\ngenetics, clinic, psychosocial, etc.)"
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Calibri"
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER


def build_presentation(structure: ParsedEssay) -> Presentation:
    """Create a PowerPoint presentation from the parsed essay structure."""
    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    title_block = structure["title_block"]

    title_shape = slide.shapes.title
    _configure_title_shape(title_shape, title_block["title"])

    subtitle_shape = slide.placeholders[1]
    subtitle_text = "\n".join(
        [
            title_block["prepared_for"],
            title_block["student"],
            title_block["course"],
        ]
    )
    subtitle_shape.text = subtitle_text
    for paragraph in subtitle_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    footer_frame = footer.text_frame
    footer_frame.text = title_block["date"]
    for run in footer_frame.paragraphs[0].runs:
        run.font.name = "Calibri"
        run.font.size = Pt(14)
    footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Section slides
    content_layout = prs.slide_layouts[1]
    for section in structure["sections"]:
        slide = prs.slides.add_slide(content_layout)
        title_shape = slide.shapes.title
        _configure_title_shape(title_shape, section["heading"], font_size=34)

        body_shape = slide.placeholders[1]
        body_shape.left = Inches(0.7)
        body_shape.width = Inches(5.5)
        bullets = _extract_bullet_points(section["paragraphs"])
        _configure_body_placeholder(body_shape, bullets)

        _add_image_placeholder(slide)

    return prs


def main() -> None:
    structure = parse_essay(ESSAY_TEXT)
    prs = build_presentation(structure)
    prs.save("Idiopathic_Short_Stature_Presentation.pptx")


if __name__ == "__main__":
    main()
